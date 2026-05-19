"""
Generate the Mini-RAG Assistant project deck (5 slides, 16:9).
Run:  python make_deck.py
Output: Mini_RAG_Assistant_Deck.pptx in the current directory.

Requires:  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---------- Theme ----------
TITLE_COLOR  = RGBColor(0x1F, 0x4E, 0x79)   # dark navy
ACCENT_COLOR = RGBColor(0x2E, 0x86, 0xAB)   # teal
TEXT_COLOR   = RGBColor(0x33, 0x33, 0x33)   # near-black

# ---------- Setup ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)   # widescreen 16:9
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_slide(title: str, subtitle: str):
    """Create a slide with a title, subtitle, and an accent bar."""
    slide = prs.slides.add_slide(BLANK)

    # Title
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4),
                                  Inches(12.3), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = TITLE_COLOR

    # Subtitle
    st = slide.shapes.add_textbox(Inches(0.5), Inches(1.3),
                                  Inches(12.3), Inches(0.5))
    sp = st.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.italic = True
    sp.font.size = Pt(18)
    sp.font.color.rgb = ACCENT_COLOR

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(1.95),
                                 Inches(2.0), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_COLOR
    bar.line.fill.background()
    return slide


def add_bullets(slide, bullets):
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.2),
                                  Inches(11.9), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + b
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)


def add_notes(slide, note: str):
    slide.notes_slide.notes_text_frame.text = note


# ---------- Slide 1 ----------
s = add_slide(
    "Problem Understanding & Objective",
    "Answers you can trust, from documents you own",
)
add_bullets(s, [
    "Problem: LLMs hallucinate - they don't know your policies, notes, or papers, and confidently make things up.",
    "Cost: wrong answers cited as fact, no source to verify, no audit trail - unacceptable in legal, medical, finance, support.",
    "Objective: a lightweight assistant that answers using ONLY a user's own documents, with inline citations and a confidence score.",
    "Constraints: free to run, deployable in minutes, no database, no shared API keys.",
    "Target user: anyone with a document corpus - policies, manuals, research papers, class notes, contracts, resumes.",
])
add_notes(s, "RAG is the industry-standard answer to LLM hallucination. "
             "The goal of this project is to build a minimal but complete "
             "implementation of the pattern, end-to-end.")


# ---------- Slide 2 ----------
s = add_slide(
    "Solution Architecture & Design Flow",
    "A 6-stage pipeline, ~290 lines of core Python",
)
add_bullets(s, [
    "Ingest - pypdf reads files, splits into ~500-word chunks with 50-word overlap.",
    "Embed - sentence-transformers/all-MiniLM-L6-v2 (local, free, ~80 MB).",
    "Index - FAISS in-memory inner-product index (cosine similarity via L2 normalization).",
    "Retrieve - top-k chunks by similarity.",
    "Generate - Llama 3.3 70B on Groq, with a strict 'cite or refuse' prompt.",
    "Score - confidence = mean similarity of the chunks the model actually cited.",
])
add_notes(s, "Nothing in this stack costs money. Each user supplies their own "
             "Groq key in the browser, so deployment requires zero secrets management.")


# ---------- Slide 3 ----------
s = add_slide(
    "Implementation Highlights",
    "Small surface area, deliberate design choices",
)
add_bullets(s, [
    "Citations parsed via regex from the answer text - no JSON mode, no function calling, no special LLM output format.",
    "Confidence grounded in retrieval, not self-assessment - LLMs are bad at knowing when they're wrong.",
    "Explicit refusal handling - one line in the prompt drops hallucinations more than any other change.",
    "L2-normalized vectors - inner product becomes cosine similarity, bounded [0, 1] scores for free.",
    "Streamlit UI built for first-time users - educational expanders, sample-question buttons, live k-mode comparison.",
    "BYO API key - no shared secrets; trivial to deploy on Hugging Face Spaces or Streamlit Cloud.",
])
add_notes(s, "The principle throughout was 'prefer simple, transparent mechanisms over clever ones'. "
             "Every design choice can be explained to a non-ML stakeholder in one sentence.")


# ---------- Slide 4 ----------
s = add_slide(
    "Challenges & Learnings",
    "What broke, what we learned",
)
add_bullets(s, [
    "Challenge: hallucination when docs lack the answer -> solved with an explicit refusal instruction in the prompt.",
    "Challenge: LLM self-reported confidence is unreliable -> replaced with retrieval-similarity-based confidence.",
    "Challenge: chunk-size trade-off -> settled on 500 words + 50 overlap as the empirical sweet spot.",
    "Challenge: scanned PDFs return empty text (pypdf can't OCR) -> flagged in UI; OCR on the roadmap.",
    "Learning: a clear, strict prompt beats a complex chain in 90% of cases.",
    "Learning: retrieval quality matters more than model size - and UX is part of the system.",
])
add_notes(s, "Biggest takeaway: most 'RAG accuracy problems' are actually retrieval or "
             "prompt problems - not LLM problems. Pick a small, fast model and invest "
             "in the surrounding pipeline.")


# ---------- Slide 5 ----------
s = add_slide(
    "Demo Summary & Next Steps",
    "What you'll see today + where this goes",
)
add_bullets(s, [
    "Demo: upload 3 sample docs -> click a sample question -> cited answer + High confidence badge.",
    "Demo: expand 'Retrieved sources' -> show the exact passages the model used.",
    "Demo: ask an out-of-scope question -> refusal triggers, confidence drops to 0.00.",
    "Next: multi-turn conversation (follow-up question rewriting) - most-requested UX gap.",
    "Next: persist FAISS index to disk + hybrid (BM25 + dense) search + cross-encoder re-ranker.",
    "Next: OCR for scanned PDFs, .docx/.xlsx/.pptx support, precision@k evaluation suite.",
])
add_notes(s, "The demo proves the core RAG loop works. The roadmap items are all "
             "incremental improvements on top of a stable foundation - nothing requires a rewrite.")


# ---------- Save ----------
out = "Mini_RAG_Assistant_Deck.pptx"
prs.save(out)
print(f"Saved {out}")
